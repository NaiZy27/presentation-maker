import asyncio
import io
import re
from pathlib import Path

from pptx import Presentation

from core.models import PresentationContent
from core.logger import get_logger
from services.images import search_image

logger = get_logger(__name__)

_TEMPLATE_10 = Path(__file__).parent.parent / "templates" / "prs_exmp_10.pptx"
_TEMPLATE_15 = Path(__file__).parent.parent / "templates" / "prs_exmp_15.pptx"


def _sanitize(text: str) -> str:
    if not isinstance(text, str):
        return str(text) if text is not None else ""
    text = re.sub(r'-{2,}', '-', text)              # -- или --- → один дефис
    text = re.sub(r'\s*[—–]\s*', ' - ', text)      # em/en dash → дефис
    for char in ('​', '﻿', '‌', '‍', '‎', '‏', ' ', ' '):
        text = text.replace(char, "")
    text = text.replace('\xa0', ' ')
    text = re.sub(r' +', ' ', text)
    text = text.replace("**", "").replace("__", "")
    # Strip leading dash/bullet markers from lines (template already has bullet formatting)
    text = re.sub(r'(?m)^[\-–•]\s+', '', text)
    return text.strip()


def _truncate_at_sentence(text: str, limit: int) -> str:
    if len(text) <= limit:
        return text
    cut = text[:limit]
    last = max(cut.rfind('.'), cut.rfind('!'), cut.rfind('?'))
    return cut[:last + 1] if last > limit // 2 else cut


def _replace_text_markers(slide, n: int, header: str, text: str) -> None:
    header_marker = f"{{{{header{n}}}}}"
    text_marker = f"{{{{text{n}}}}}"

    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                if header_marker in run.text:
                    run.text = run.text.replace(header_marker, header)
                if text_marker in run.text:
                    run.text = run.text.replace(text_marker, text)


def _replace_image_group(slide, n: int, image_bytes: bytes) -> None:
    image_marker = f"{{{{image{n}}}}}"

    for shape in slide.shapes:
        if shape.name == image_marker:
            left, top, width, height = shape.left, shape.top, shape.width, shape.height
            slide.shapes._spTree.remove(shape._element)
            slide.shapes.add_picture(io.BytesIO(image_bytes), left, top, width=width, height=height)
            logger.info("Image inserted at slide %d", n)
            return

    logger.warning("Image group %r not found on slide %d", image_marker, n)


def _build_sync(content: PresentationContent, images: list[bytes | None], output_path: str, slide_count: int) -> None:
    template = _TEMPLATE_15 if slide_count == 15 else _TEMPLATE_10
    prs = Presentation(str(template))

    for i, slide_content in enumerate(content.slides):
        n = i + 1
        slide = prs.slides[i + 1]  # slide index 0 is title, content starts at 1

        _replace_text_markers(
            slide, n,
            header=_sanitize(slide_content.header)[:80],
            text=_truncate_at_sentence(_sanitize(slide_content.text), 1200),
        )

        if images[i] is not None:
            _replace_image_group(slide, n, images[i])
        else:
            logger.warning("No image for slide %d, placeholder left", n)

    prs.save(output_path)
    logger.info("Presentation saved to %s", output_path)


async def build_presentation(content: PresentationContent, output_path: str, slide_count: int = 10) -> None:
    logger.info("Downloading %d images...", len(content.slides))
    images = await asyncio.gather(
        *[search_image(slide.image_query) for slide in content.slides]
    )

    logger.info("Building pptx...")
    await asyncio.to_thread(_build_sync, content, list(images), output_path, slide_count)
